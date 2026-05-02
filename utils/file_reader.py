"""
utils/file_reader.py — Membaca isi file PDF, Word, dan PowerPoint
"""
import os
from pathlib import Path


def read_pdf(file_path: str) -> str:
    """Membaca teks dari file PDF menggunakan PyMuPDF"""
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(file_path)
        text = ""
        for page in doc:
            text += page.get_text() + "\n"
        doc.close()
        return text.strip()
    except Exception as e:
        print(f"  ❌ Gagal baca PDF: {e}")
        return ""


def read_docx(file_path: str) -> str:
    """Membaca teks dari file Word (.docx)"""
    try:
        from docx import Document
        doc = Document(file_path)
        text = ""
        for para in doc.paragraphs:
            text += para.text + "\n"
        return text.strip()
    except Exception as e:
        print(f"  ❌ Gagal baca DOCX: {e}")
        return ""


def read_pptx(file_path: str) -> str:
    """Membaca teks dari file PowerPoint (.pptx)"""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        text = ""
        for slide_num, slide in enumerate(prs.slides, 1):
            text += f"\n--- Slide {slide_num} ---\n"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text += shape.text + "\n"
        return text.strip()
    except Exception as e:
        print(f"  ❌ Gagal baca PPTX: {e}")
        return ""


def read_file(file_path: str) -> str:
    """
    Membaca isi file berdasarkan ekstensi.
    Mendukung: .pdf, .docx, .doc, .pptx, .ppt, .txt
    """
    ext = Path(file_path).suffix.lower()
    
    if ext == ".pdf":
        return read_pdf(file_path)
    elif ext in (".docx", ".doc"):
        return read_docx(file_path)
    elif ext in (".pptx", ".ppt"):
        return read_pptx(file_path)
    elif ext == ".txt":
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                return f.read()
        except Exception as e:
            print(f"  ❌ Gagal baca TXT: {e}")
            return ""
    else:
        print(f"  ⚠️ Format file {ext} tidak didukung untuk dibaca")
        return ""
